"""Build Kotak Agentic Reader presentation matching Transaction_Tagger_Presentation.pptx style."""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ----- Palette -----
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x27, 0x5B, 0xAD)
GOLD   = RGBColor(0xB8, 0x86, 0x0B)
TINT   = RGBColor(0xE8, 0xED, 0xF7)
BORDER = RGBColor(0xCC, 0xD3, 0xDD)
TEXT   = RGBColor(0x55, 0x65, 0x78)
SUBTLE = RGBColor(0x99, 0xA8, 0xBB)
GIANT  = RGBColor(0x22, 0x4A, 0x85)
LBLUE  = RGBColor(0x88, 0xA8, 0xD8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xF1, 0x9A)

prs = Presentation()
prs.slide_width  = Emu(12188952)
prs.slide_height = Emu(6858000)

blank = prs.slide_layouts[6]
TOTAL = 13


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_chrome(slide, idx, title, subtitle, footer):
    """Top + left strips, slide number, title, divider, subtitle, footer bar."""
    # left vertical navy strip
    add_rect(slide, Emu(0), Emu(0), Emu(64008), Emu(6858000), NAVY)
    # top horizontal navy strip
    add_rect(slide, Emu(0), Emu(0), Emu(12188952), Emu(36576), NAVY)
    # slide number top-right
    add_text(slide, Emu(11365992), Emu(109728), Emu(685800), Emu(320040),
             f"{idx:02d} / {TOTAL:02d}", size=10, color=TEXT, align=PP_ALIGN.RIGHT)
    # title
    add_text(slide, Emu(274320), Emu(164592), Emu(10515600), Emu(502920),
             title, size=24, bold=True, color=NAVY)
    # divider under title
    add_rect(slide, Emu(274320), Emu(713232), Emu(11612880), Emu(15240), BORDER)
    # subtitle
    add_text(slide, Emu(274320), Emu(768096), Emu(10058400), Emu(365760),
             subtitle, size=12, color=TEXT)
    # footer bar (light tint)
    add_rect(slide, Emu(0), Emu(6510528), Emu(12188952), Emu(347472), TINT)
    add_rect(slide, Emu(0), Emu(6510528), Emu(12188952), Emu(12700), BORDER)
    add_text(slide, Emu(274320), Emu(6547104), Emu(11612880), Emu(274320),
             footer, size=8, color=TEXT)


# ============================================================
# SLIDE 1 — Title / cover
# ============================================================
s = prs.slides.add_slide(blank)
# left dark panel
add_rect(s, Emu(0), Emu(0), Emu(4206240), Emu(6858000), NAVY)
add_rect(s, Emu(0), Emu(5029200), Emu(4206240), Emu(54864), BLUE)
add_rect(s, Emu(0), Emu(5120640), Emu(4206240), Emu(27432), GOLD)

# giant "RG" mark on the dark panel
add_text(s, Emu(91440), Emu(914400), Emu(4023360), Emu(3657600),
         "RG", size=220, bold=True, color=GIANT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# top kicker
add_text(s, Emu(228600), Emu(201168), Emu(3749039), Emu(274320),
         "CONFIDENTIAL  ·  AGENTIC AI BRIEFING", size=8, color=LBLUE)

# main title (right side)
add_text(s, Emu(4572000), Emu(1097280), Emu(7132320), Emu(2560320),
         "Kotak Agentic\nReader\nfor Hindsighting", size=40, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# gold divider
add_rect(s, Emu(4572000), Emu(3703320), Emu(2743200), Emu(38100), GOLD)

# subtitle
add_text(s, Emu(4572000), Emu(3840480), Emu(7132320), Emu(822960),
         "Streamlining root-cause analysis for offers, first payment defaults,\nand 30+ DPD cases — from manual case-by-case review to an agent-assisted workflow.",
         size=13, color=TEXT)

# stat cards (4)
def stat_card(x, num, label):
    add_rect(s, x, Emu(4937760), Emu(1645920), Emu(749808), TINT, line=BORDER)
    add_rect(s, x, Emu(4937760), Emu(45720), Emu(749808), NAVY)
    add_text(s, x + Emu(91440), Emu(4992624), Emu(1554480), Emu(360000),
             num, size=20, bold=True, color=NAVY)
    add_text(s, x + Emu(91440), Emu(5380000), Emu(1554480), Emu(280000),
             label, size=8, color=TEXT)

stat_card(Emu(4572000),  "10–20",   "Cases / day today (manual)")
stat_card(Emu(6382512),  "3",       "Failure modes diagnosed")
stat_card(Emu(8193024),  "100%",    "Local LLM · zero PII egress")
stat_card(Emu(10003536), "HTML+XLSX", "Reports shipped")

# bottom byline
add_text(s, Emu(4572000), Emu(6446520), Emu(7315200), Emu(320040),
         "Built by the AI Team  ·  LangChain · Ollama · Deterministic-first design",
         size=8, color=SUBTLE)


# ============================================================
# SLIDE 2 — Context, problem & what we shipped (60/40 split)
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 2,
           "The Hindsighting Bottleneck — and What We Shipped",
           "Why this project exists, what can go wrong with an offer, and where we are today.",
           "Cases/day (manual): 10–20  ·  Failure modes covered: execution · policy · income estimation  ·  Shipped: HTML + Excel reports + agentic base")

# 60/40 split
content_top = Emu(1280000)
content_h   = Emu(4720000)
gap_x       = Emu(140000)
total_w     = Emu(11548872)
left_w  = Emu(int(total_w * 0.60)) - Emu(int(gap_x/2))
right_w = total_w - left_w - gap_x
x_left  = Emu(320040)
x_right = x_left + left_w + gap_x

# Left column = 2 stacked rows
row_gap = Emu(120000)
row_h   = Emu(int((content_h - row_gap) / 2))

def panel(x, y, w, h, header, body_lines, accent):
    add_rect(s, x, y, w, h, TINT)
    add_rect(s, x, y, w, Emu(60000), accent)
    add_text(s, x+Emu(180000), y+Emu(130000), w-Emu(360000), Emu(400000),
             header, size=13, bold=True, color=NAVY)
    add_text(s, x+Emu(180000), y+Emu(610000), w-Emu(360000), h-Emu(700000),
             "\n\n".join(body_lines), size=11, color=TEXT)

# Top-left: Where we are today
panel(x_left, content_top, left_w, row_h,
      "Where we are today",
      ["•  Analysts manually pull data, aggregate, and review 10–20 cases / day.",
       "•  Each case asks the same question: what went wrong with this offer / FPD / 30+ DPD?",
       "•  Loop is slow, inconsistent across analysts, and hard to audit."],
      accent=NAVY)

# Bottom-left: What we shipped
panel(x_left, content_top + row_h + row_gap, left_w, row_h,
      "What we shipped",
      ["•  Customer + Bureau + Combined HTML / PDF / Excel reports.",
       "•  Transaction tagging in-built for Kotak banking narrations (43-category model wired in as a tool).",
       "•  Deterministic engine: feature vectors, key findings, scorecard, checklist.",
       "•  Agentic base: intent parser, planner, tool registry, audit log — ready to evolve into a ReAct loop."],
      accent=BLUE)

# Right column: Three failure modes (full height)
panel(x_right, content_top, right_w, content_h,
      "Three failure modes we diagnose",
      ["•  Execution gaps\n    a policy or sub-component did not fire correctly.",
       "•  Policy issues\n    too harsh, or safe to relax based on observed cohort behaviour.",
       "•  Income estimation\n    under- or over-stated; alternate income signals were missed."],
      accent=GOLD)


# ============================================================
# SLIDE 3 — Generic architecture
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 3,
           "System Architecture",
           "Five-stage agentic pipeline, deterministic tool layer, and multi-format outputs — all running on local Ollama.",
           "Mistral (intent · JSON)  ·  llama3.2 (narration)  ·  pandas + fpdf2 + Jinja2  ·  20+ tools registered today")
s.shapes.add_picture("/tmp/ppt_assets/arch_generic.png",
                     Emu(320040), Emu(1230000), width=Emu(11548000))


# ============================================================
# SLIDE 4 — Determinism > Intelligence
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 4,
           "Determinism > Intelligence",
           "Every number on every report is computed in pure Python. The LLM never sees raw data and never decides risk.",
           "Why: regulatory auditability  ·  reproducibility  ·  no hallucinated numbers  ·  single source of truth in config/thresholds.py")
s.shapes.add_picture("/tmp/ppt_assets/arch_determinism.png",
                     Emu(320040), Emu(1230000), width=Emu(9100000))

# Right column: motivation — concise bullets only
right_x = Emu(9560000); right_y = Emu(1230000)
right_w = Emu(12188952) - right_x - Emu(320040)
panel_h = Emu(5180000)

add_rect(s, right_x, right_y, right_w, panel_h, TINT)
add_rect(s, right_x, right_y, right_w, Emu(60000), GOLD)

add_text(s, right_x+Emu(160000), right_y+Emu(130000), right_w-Emu(320000), Emu(280000),
         "MOTIVATION", size=8, bold=True, color=GOLD)

bullets = [
    "Auditable line-by-line",
    "Same input → same output",
    "LLM never sees raw data",
    "Narration fails soft — numbers always render",
    "Findings = constant + comparison op",
]
yy = right_y + Emu(500000)
for b in bullets:
    add_text(s, right_x+Emu(160000), yy, right_w-Emu(320000), Emu(360000),
             "•  " + b, size=11, color=NAVY, bold=True)
    yy += Emu(420000)


# ============================================================
# SLIDE 5 — Feature extraction + Key Findings
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 5,
           "Bureau Feature Extraction & Key-Findings Engine",
           "100+ raw loan types normalized to 13 canonical classes; per-type feature vectors feed a 40+ rule, zero-LLM findings engine.",
           "Severity ladder: high_risk · moderate_risk · concern · positive · neutral  ·  Same constants annotate the LLM narration prompt")
s.shapes.add_picture("/tmp/ppt_assets/arch_features.png",
                     Emu(320040), Emu(1230000), width=Emu(11548000))


# ============================================================
# SLIDE 6 — Transaction Tagger (upstream) — HYPERLINK
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 6,
           "Upstream Dependency — Transaction Tagger",
           "Banking narrations are first classified into 43 categories by a separate retrieval-augmented model. That work is documented in its own deck.",
           "Click the panel to open the dedicated Transaction Tagger presentation  ·  Adapted from Amazon Science research")

# Big clickable card
card_x = Emu(1500000); card_y = Emu(1500000); card_w = Emu(9200000); card_h = Emu(4400000)
card = add_rect(s, card_x, card_y, card_w, card_h, NAVY)
# gold accent stripe
add_rect(s, card_x, card_y, card_w, Emu(60000), GOLD)
add_rect(s, card_x, card_y, Emu(60000), card_h, GOLD)

add_text(s, card_x+Emu(400000), card_y+Emu(400000), card_w-Emu(800000), Emu(600000),
         "TRANSACTION TAGGER", size=10, bold=True, color=GOLD)

add_text(s, card_x+Emu(400000), card_y+Emu(900000), card_w-Emu(800000), Emu(900000),
         "Intelligence-Led Transaction Classification",
         size=28, bold=True, color=WHITE)

add_text(s, card_x+Emu(400000), card_y+Emu(2000000), card_w-Emu(800000), Emu(1200000),
         "Retrieval-augmented categorization of banking narrations.\n"
         "Powers the category_of_txn column that every downstream feature, finding,\n"
         "and report section depends on.",
         size=14, color=LBLUE)

# Stat strip inside card
strip_y = card_y + Emu(3200000)
def mini_stat(x, num, label):
    add_text(s, x, strip_y, Emu(2000000), Emu(400000), num, size=20, bold=True, color=GOLD)
    add_text(s, x, strip_y+Emu(420000), Emu(2000000), Emu(280000), label, size=9, color=LBLUE)
mini_stat(card_x+Emu(400000),  "43",     "Categories")
mini_stat(card_x+Emu(2700000), "35K+",   "Labeled records")
mini_stat(card_x+Emu(5000000), "10M+",   "Txns / month")
mini_stat(card_x+Emu(7300000), "<15ms",  "Per transaction")

# Hyperlink the whole card to the tagger PPTX
def add_hyperlink(shape, target):
    # add a click action with hyperlink
    spPr = shape.click_action
    spPr.hyperlink.address = target

add_hyperlink(card, "Transaction_Tagger_Presentation.pptx")

# Call to action below
add_text(s, card_x, card_y+card_h+Emu(150000), card_w, Emu(300000),
         "▶  Click the panel above to open the full Transaction Tagger deck",
         size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — Why this pipeline needs an agent
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 7,
           "Why This Pipeline Needs an Agent — Report Orchestration",
           "3 builders × 5 themes × 3 output formats × per-section narration chains × caching. The combinatorics is what makes agentic orchestration the right paradigm — not a luxury.",
           "Today: deterministic routing through INTENT_TOOL_MAP  ·  Tomorrow: full ReAct loop (see next slide)")

# Diagram — bigger, takes most of the slide
img_x = Emu(320040); img_y = Emu(1230000); img_w = Emu(9100000)
s.shapes.add_picture("/tmp/ppt_assets/arch_rendering.png", img_x, img_y, width=img_w)

# Right panel: concise "Why agentic" + scaling message
panel_x = img_x + img_w + Emu(180000)
panel_y = img_y
panel_w = Emu(12188952) - panel_x - Emu(320040)
panel_h = Emu(5180000)

add_rect(s, panel_x, panel_y, panel_w, panel_h, TINT)
add_rect(s, panel_x, panel_y, panel_w, Emu(60000), GOLD)

add_text(s, panel_x+Emu(160000), panel_y+Emu(130000), panel_w-Emu(320000), Emu(280000),
         "WHY AGENTIC?", size=8, bold=True, color=GOLD)

bullets = [
    "Pick the right builder",
    "Skip empty sections",
    "Reuse cache or rebuild",
    "Retry on LLM failure",
    "Choose the output format",
]
yy = panel_y + Emu(500000)
for b in bullets:
    add_text(s, panel_x+Emu(160000), yy, panel_w-Emu(320000), Emu(320000),
             "•  " + b, size=11, color=NAVY, bold=True)
    yy += Emu(380000)

# Scaling tagline (the headline message)
tag_y = panel_y + panel_h - Emu(1700000)
add_rect(s, panel_x+Emu(140000), tag_y, panel_w-Emu(280000), Emu(15240), GOLD)
add_text(s, panel_x+Emu(160000), tag_y+Emu(80000), panel_w-Emu(320000), Emu(1500000),
         "More agents.\nMore data sources.\nMore decisions per case.\n\nThe agent scales linearly — manual review does not.",
         size=11, bold=True, color=NAVY)


# ============================================================
# SLIDE 8 — Agentic foundation → ReAct
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 8,
           "Agentic Foundation — Toward a ReAct Agent",
           "Today's pipeline is the deterministic skeleton of a ReAct agent. The remaining step is to replace the static INTENT_TOOL_MAP with dynamic LLM-driven tool selection and a reflection loop.",
           "Already built: parser · planner · tool registry · audit log  ·  Gap: dynamic tool selection + observation-driven re-planning")
s.shapes.add_picture("/tmp/ppt_assets/arch_agentic.png",
                     Emu(320040), Emu(1230000), width=Emu(11548000))


# ============================================================
# SLIDE 9 — Sample output: Combined Report
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 9,
           "Sample Output — Combined Report (HTML)",
           "Real customer view: scorecard, bureau key findings, checklist, narrative — generated end-to-end in seconds.",
           "Sample CRN: 698167220  ·  Yellow highlights to be added per analyst walkthrough")

# Placeholder box
ph_x = Emu(700000); ph_y = Emu(1300000); ph_w = Emu(10800000); ph_h = Emu(4900000)
add_rect(s, ph_x, ph_y, ph_w, ph_h, WHITE, line=BORDER)
add_rect(s, ph_x, ph_y, ph_w, Emu(60000), GOLD)
add_text(s, ph_x, ph_y+Emu(2100000), ph_w, Emu(500000),
         "[ Insert annotated Combined Report screenshot ]",
         size=18, bold=True, color=SUBTLE, align=PP_ALIGN.CENTER)
add_text(s, ph_x, ph_y+Emu(2700000), ph_w, Emu(400000),
         "Yellow highlights mark the analyst's review points",
         size=11, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, ph_x, ph_y+Emu(3100000), ph_w, Emu(300000),
         "reports/combined_698167220_report.html",
         size=9, color=SUBTLE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 — Sample output: Excel / Checklist / Scorecard
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 10,
           "Sample Output — Excel Workbook · Checklist · Scorecard",
           "Structured outputs for downstream policy review, batch QC, and audit trail.",
           "Excel sheet per section  ·  Checklist booleans drive policy diff  ·  Scorecard is pure deterministic")

# Two side-by-side placeholders
def half_ph(x, label):
    ph_w = Emu(5600000); ph_h = Emu(4900000); ph_y = Emu(1300000)
    add_rect(s, x, ph_y, ph_w, ph_h, WHITE, line=BORDER)
    add_rect(s, x, ph_y, ph_w, Emu(60000), GOLD)
    add_text(s, x, ph_y+Emu(2100000), ph_w, Emu(500000),
             label, size=15, bold=True, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_text(s, x, ph_y+Emu(2650000), ph_w, Emu(400000),
             "[ screenshot to be inserted ]",
             size=10, color=TEXT, align=PP_ALIGN.CENTER)

half_ph(Emu(320040),  "Excel Workbook — multi-sheet")
half_ph(Emu(6240000), "Checklist + Scorecard view")


# ============================================================
# SLIDE 11 — Impact (directional)
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 11,
           "Impact — Directional",
           "We do not yet have hard TAT numbers; the qualitative shift is the headline.",
           "Pilot phase: shipping HTML + Excel to analysts for shadow-mode comparison against manual reviews")

# 4 stacked rows (same panel design as roadmap, simpler content)
points = [
    ("THROUGHPUT",   "Case TAT compression",
     "Manual data pull + aggregation collapses to a single agent invocation. Analyst time shifts from gathering to judging.", BLUE),
    ("CONSISTENCY",  "Same finding every time",
     "Deterministic thresholds eliminate analyst-to-analyst variance on the same input cohort.", GOLD),
    ("AUDITABILITY", "Every number traceable",
     "Each key finding maps to a named constant in config/thresholds.py and a single comparison operator.", NAVY),
    ("SCALABILITY",  "From 10–20 cases / day → batch",
     "Same pipeline runs interactively or via batch_reports.py over arbitrary CRN lists.", BLUE),
]
y = Emu(1280000); h = Emu(1100000); gap = Emu(120000); w = Emu(11540000); x = Emu(320040)
for i, (kicker, head, body, accent) in enumerate(points):
    yy = y + i*(h+gap)
    add_rect(s, x, yy, w, h, TINT)
    add_rect(s, x, yy, Emu(60000), h, accent)
    add_text(s, x+Emu(180000), yy+Emu(140000), Emu(10000000), Emu(260000),
             kicker, size=8, bold=True, color=accent)
    add_text(s, x+Emu(180000), yy+Emu(400000), Emu(10000000), Emu(360000),
             head, size=14, bold=True, color=NAVY)
    add_text(s, x+Emu(180000), yy+Emu(760000), Emu(11000000), Emu(320000),
             body, size=11, color=TEXT)


# ============================================================
# SLIDE 12 — Roadmap
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 12,
           "Roadmap",
           "From shipped HTML/Excel reports to a fully agentic offer-validation system.",
           "Each phase builds on the prior one — no rewrite, only extension of the existing pipeline")

phases = [
    ("Phase 1 — NOW",      "Combined HTML + Excel reports. Deterministic engine. Static intent → tool routing.",        NAVY),
    ("Phase 2 — ReAct Loop", "Replace INTENT_TOOL_MAP with LLM-driven tool selection + observation-based re-planning.", BLUE),
    ("Phase 3 — Offer Validation",  "Forensic playbooks for FPD / 30+ DPD: policy diff, income re-estimation, execution-gap detector.", GOLD),
    ("Phase 4 — Active Learning",   "Analyst overrides feed back into threshold tuning and category-mapping refinement.", NAVY),
]
y = Emu(1280000); h = Emu(1100000); gap = Emu(120000); w = Emu(11540000); x = Emu(320040)
for i, (head, body, accent) in enumerate(phases):
    yy = y + i*(h+gap)
    add_rect(s, x, yy, w, h, TINT)
    add_rect(s, x, yy, Emu(60000), h, accent)
    # number badge
    add_rect(s, x+Emu(140000), yy+Emu(310000), Emu(290000), Emu(290000), accent)
    add_text(s, x+Emu(140000), yy+Emu(290000), Emu(290000), Emu(330000),
             str(i+1), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x+Emu(540000), yy+Emu(150000), Emu(10000000), Emu(360000),
             head, size=13, bold=True, color=NAVY)
    add_text(s, x+Emu(540000), yy+Emu(520000), Emu(10000000), Emu(500000),
             body, size=11, color=TEXT)


# ============================================================
# SLIDE 13 — Risks, guardrails + Ask
# ============================================================
s = prs.slides.add_slide(blank)
add_chrome(s, 13,
           "Risks, Guardrails & Ask",
           "What can break, what protects us, and what we need next.",
           "Confidential  ·  Internal review only")

# Two columns
col_w = Emu(5600000); col_h = Emu(4700000); top = Emu(1280000); gap = Emu(120000)

def colbox(x, header, items, accent):
    add_rect(s, x, top, col_w, col_h, TINT)
    add_rect(s, x, top, col_w, Emu(60000), accent)
    add_text(s, x+Emu(180000), top+Emu(150000), col_w-Emu(360000), Emu(400000),
             header, size=14, bold=True, color=NAVY)
    add_text(s, x+Emu(180000), top+Emu(670000), col_w-Emu(360000), col_h-Emu(800000),
             "\n\n".join(items), size=11, color=TEXT)

colbox(Emu(320040), "Guardrails in place",
       ["•  Local Ollama inference — zero PII egress.",
        "•  Fail-soft everywhere: tool/LLM failure never crashes the pipeline.",
        "•  JSONL audit log per query (logs/audit_*.jsonl).",
        "•  Deterministic numbers — LLM cannot mutate a finding.",
        "•  Single-source thresholds, prompts, and categories under config/."],
       accent=NAVY)

colbox(Emu(6240000), "Risks & the Ask",
       ["•  Single-user process today — needs hardening for concurrent analysts.",
        "•  Cache invalidation is manual — needs TTL for daily data refresh.",
        "•  Pilot ask: 2 analysts × 4 weeks shadow-mode comparison.",
        "•  Infra ask: dedicated Ollama host + analyst sandbox.",
        "•  Decision ask: green-light Phase 2 (ReAct) scope and timeline."],
       accent=GOLD)


# ============================================================
out = "Kotak_Agentic_Reader_Presentation.pptx"
prs.save(out)
print(f"✓ Saved {out} with {len(prs.slides)} slides")
